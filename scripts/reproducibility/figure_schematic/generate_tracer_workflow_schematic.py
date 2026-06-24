#!/usr/bin/env python
"""Generate the TRACER workflow schematic.

This script creates a PowerPoint-native, vector-first Figure 1A-style
schematic for TRACER and exports SVG/PDF/PNG companions from the same
layout. The PPTX uses editable PowerPoint shapes and editable text for
the figure body; the only raster asset is the provided TRACER logo.
"""
from __future__ import annotations

import base64
import math
from pathlib import Path

import cairosvg
import svgwrite
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT = ROOT / "outputs"
LOGO = ROOT / "assets" / "images" / "logo.png"

PPTX_OUT = OUT / "tracer_workflow_schematic.pptx"
SVG_OUT = OUT / "tracer_workflow_schematic.svg"
PDF_OUT = OUT / "tracer_workflow_schematic.pdf"
PNG_OUT = OUT / "tracer_workflow_schematic.png"

W, H = 1600, 900
SLIDE_W, SLIDE_H = 13.333333, 7.5

COL = {
    "bg": "F7FAFC",
    "paper": "FFFFFF",
    "ink": "172A3A",
    "muted": "5B6775",
    "line": "C9D5DF",
    "line2": "DDE7EE",
    "cyan": "2CA9B7",
    "cyan_dark": "0E7480",
    "teal": "2EA66F",
    "green": "7BBE5B",
    "indigo": "5867D8",
    "violet": "8666C6",
    "coral": "E36F5E",
    "amber": "D9A441",
    "slate": "34495E",
    "soft_cyan": "E8F7FA",
    "soft_teal": "EAF8F0",
    "soft_indigo": "EEF0FE",
    "soft_coral": "FDEDEB",
    "soft_amber": "FFF6DE",
    "soft_slate": "EEF3F7",
}


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.replace("#", "")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def px_to_in(v: float) -> float:
    return v / W * SLIDE_W


def py_to_in(v: float) -> float:
    return v / H * SLIDE_H


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: str = "ink",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = "Aptos",
    line_spacing: float = 0.95,
):
    box = slide.shapes.add_textbox(Inches(px_to_in(x)), Inches(py_to_in(y)), Inches(px_to_in(w)), Inches(py_to_in(h)))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COL[color])
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = "paper",
    line: str = "line",
    radius=True,
    width: float = 1.0,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(px_to_in(x)), Inches(py_to_in(y)), Inches(px_to_in(w)), Inches(py_to_in(h)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COL[fill])
    shp.line.color.rgb = rgb(COL[line])
    shp.line.width = Pt(width)
    return shp


def add_circle(slide, cx: float, cy: float, r: float, *, fill: str, line: str = "paper", width: float = 1.0):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(px_to_in(cx - r)),
        Inches(py_to_in(cy - r)),
        Inches(px_to_in(2 * r)),
        Inches(py_to_in(2 * r)),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COL[fill])
    shp.line.color.rgb = rgb(COL[line])
    shp.line.width = Pt(width)
    return shp


def add_line(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = "line", width: float = 1.2):
    line = slide.shapes.add_connector(
        1,
        Inches(px_to_in(x1)),
        Inches(py_to_in(y1)),
        Inches(px_to_in(x2)),
        Inches(py_to_in(y2)),
    )
    line.line.color.rgb = rgb(COL[color])
    line.line.width = Pt(width)
    return line


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = "cyan", width: float = 2.0):
    add_line(slide, x1, y1, x2, y2, color=color, width=width)
    ang = math.atan2(y2 - y1, x2 - x1)
    size = 13
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        Inches(px_to_in(x2 - size / 2)),
        Inches(py_to_in(y2 - size / 2)),
        Inches(px_to_in(size)),
        Inches(py_to_in(size)),
    )
    shp.rotation = math.degrees(ang) + 90
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COL[color])
    shp.line.color.rgb = rgb(COL[color])
    return shp


def add_label_pill(slide, text: str, x: float, y: float, w: float, *, fill: str, color: str = "ink", size: float = 9.0):
    add_rect(slide, x, y, w, 28, fill=fill, line=fill, radius=True, width=0)
    add_text(slide, text, x + 9, y + 6, w - 18, 16, size=size, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_card(slide, title: str, subtitle: str, x: float, y: float, w: float, h: float, *, accent: str, fill: str):
    add_rect(slide, x + 4, y + 6, w, h, fill="line2", line="line2", radius=True, width=0)
    add_rect(slide, x, y, w, h, fill=fill, line="line", radius=True, width=0.8)
    add_rect(slide, x, y, 6, h, fill=accent, line=accent, radius=False, width=0)
    add_text(slide, title, x + 18, y + 16, w - 36, 22, size=13.5, color="ink", bold=True)
    add_text(slide, subtitle, x + 18, y + 42, w - 36, h - 50, size=8.4, color="muted", line_spacing=0.92)


def add_micro_cell(slide, cx: float, cy: float, scale: float, *, fill: str = "soft_cyan", nucleus: str = "indigo"):
    add_circle(slide, cx, cy, 26 * scale, fill=fill, line="cyan", width=0.9)
    add_circle(slide, cx - 4 * scale, cy - 2 * scale, 9 * scale, fill=nucleus, line="paper", width=0.6)
    for dx, dy, c in [(-18, -12, "coral"), (12, -15, "teal"), (19, 8, "amber"), (-10, 15, "indigo")]:
        add_circle(slide, cx + dx * scale, cy + dy * scale, 3.2 * scale, fill=c, line=c, width=0)


def add_bin_grid(slide, x: float, y: float, s: float):
    for i in range(4):
        for j in range(4):
            fill = ["soft_cyan", "soft_indigo", "soft_teal", "soft_amber"][(i + j) % 4]
            add_rect(slide, x + j * s, y + i * s, s - 2, s - 2, fill=fill, line="line2", radius=False, width=0.5)
    for px, py, c in [(1.5, 1.2, "coral"), (2.4, 2.5, "teal"), (0.7, 2.2, "indigo")]:
        add_circle(slide, x + px * s, y + py * s, 5, fill=c, line=c, width=0)


def add_tissue_problem(slide):
    x, y = 405, 168
    add_rect(slide, x, y, 310, 132, fill="paper", line="line", radius=True, width=0.8)
    add_text(slide, "Why segmentation breaks", x + 18, y + 13, 210, 20, size=12.5, bold=True)
    add_text(slide, "3D tissue volume projected to 2D", x + 18, y + 36, 178, 17, size=8.2, color="muted")
    for k, (dx, dy, col) in enumerate([(0, 0, "soft_indigo"), (12, 12, "soft_cyan"), (24, 24, "soft_teal")]):
        add_rect(slide, x + 25 + dx, y + 66 + dy, 108, 42, fill=col, line="line2", radius=True, width=0.6)
        add_micro_cell(slide, x + 62 + dx, y + 86 + dy, 0.55, fill="paper", nucleus="indigo")
        add_micro_cell(slide, x + 98 + dx, y + 90 + dy, 0.50, fill="paper", nucleus="coral")
    add_line(slide, x + 150, y + 62, x + 206, y + 100, color="coral", width=1.4)
    add_line(slide, x + 150, y + 118, x + 208, y + 90, color="coral", width=1.4)
    add_text(slide, "overlap\nmixed genes\npartial / anuclear", x + 198, y + 69, 92, 56, size=8.1, color="muted", line_spacing=0.88)


def add_pptx():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, fill="bg", line="bg", radius=False, width=0)

    # Header
    slide.shapes.add_picture(str(LOGO), Inches(px_to_in(44)), Inches(py_to_in(35)), width=Inches(px_to_in(142)))
    add_text(slide, "TRACER", 206, 34, 145, 38, size=24, bold=True, color="ink")
    add_text(slide, "Transcript Coherence-Aware Reconstruction for spatial transcriptomics", 207, 71, 510, 24, size=10.5, color="muted")
    add_text(slide, "From noisy 2D transcript assignments to coherent whole-cell and partial-cell profiles", 815, 42, 640, 26, size=13.5, color="ink", bold=True, align=PP_ALIGN.RIGHT)
    add_line(slide, 44, 113, 1548, 113, color="line2", width=0.9)

    # Input modalities
    add_card(slide, "Imaging-based ST", "Segmented or transcript-level input\nXenium | Xenium5K | Atera\nCosMx | MERFISH", 52, 170, 280, 170, accent="cyan", fill="paper")
    add_micro_cell(slide, 100, 292, 0.78, fill="soft_cyan", nucleus="indigo")
    add_micro_cell(slide, 148, 286, 0.66, fill="soft_teal", nucleus="teal")
    for i, c in enumerate(["coral", "amber", "indigo", "teal", "cyan"]):
        add_circle(slide, 215 + i * 14, 303 - (i % 2) * 10, 4, fill=c, line=c, width=0)

    add_card(slide, "Sequencing-based ST", "VisiumHD, bin-level or pixel-level\nNo prior segmentation required", 52, 370, 280, 155, accent="teal", fill="paper")
    add_bin_grid(slide, 82, 448, 18)
    add_text(slide, "ambiguous\nbins / pixels", 178, 457, 102, 44, size=8.1, color="muted", line_spacing=0.9)

    add_arrow(slide, 335, 276, 390, 276, color="cyan", width=2.2)
    add_arrow(slide, 335, 448, 390, 448, color="teal", width=2.2)
    add_tissue_problem(slide)
    add_arrow(slide, 720, 245, 770, 245, color="slate", width=1.8)

    # Method core container
    add_rect(slide, 760, 148, 445, 222, fill="paper", line="line", radius=True, width=0.9)
    add_label_pill(slide, "TRACER method core", 784, 165, 150, fill="soft_indigo", color="indigo")
    steps = [
        ("1", "Coherence\nmap", "gene-pair\nevidence", "indigo", "soft_indigo"),
        ("2", "Conflict\nprune", "incompatible\nassignments", "coral", "soft_coral"),
        ("3", "Rescue &\nregroup", "spatial\nwitnesses", "teal", "soft_teal"),
        ("4", "Stitch /\nrefine", "coherent\nfragments", "cyan", "soft_cyan"),
        ("5", "Partial-cell\nreconstruct", "anuclear /\ncut cells", "amber", "soft_amber"),
    ]
    sx = 782
    for idx, (num, title, sub, accent, fill) in enumerate(steps):
        x = sx + idx * 82
        add_circle(slide, x + 30, 216, 23, fill=fill, line=accent, width=1.1)
        add_text(slide, num, x + 18, 203, 24, 20, size=12, bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_text(slide, title, x - 2, 247, 70, 36, size=9.3, bold=True, align=PP_ALIGN.CENTER, line_spacing=0.85)
        add_text(slide, sub, x - 9, 289, 88, 45, size=6.6, color="muted", align=PP_ALIGN.CENTER, line_spacing=0.88)
        if idx < len(steps) - 1:
            add_arrow(slide, x + 57, 216, x + 78, 216, color="line", width=1.2)

    # Mini mechanism sketches
    add_rect(slide, 792, 333, 378, 22, fill="soft_slate", line="soft_slate", radius=True, width=0)
    add_text(slide, "composition score C(entity) + spatial gates + witness-supported reassignment", 808, 338, 344, 12, size=7.1, color="slate", align=PP_ALIGN.CENTER)

    # Operation modes
    add_rect(slide, 405, 595, 558, 190, fill="paper", line="line", radius=True, width=0.9)
    add_text(slide, "Operation modes", 430, 617, 160, 20, size=13, bold=True)
    modes = [
        ("Refine in place", "preserve original whole cells\nprune conflict; reconstruct partials", "cyan", "soft_cyan"),
        ("Resegment", "start from segmentation\nallow new entities and stitching", "indigo", "soft_indigo"),
        ("Noseg", "discard labels\nreconstruct from bins / pixels", "teal", "soft_teal"),
    ]
    for i, (title, body, accent, fill) in enumerate(modes):
        x = 430 + i * 172
        add_rect(slide, x, 653, 150, 96, fill=fill, line=accent, radius=True, width=0.8)
        add_circle(slide, x + 27, 681, 14, fill=accent, line=accent, width=0)
        if i == 0:
            add_micro_cell(slide, x + 27, 681, 0.30, fill="paper", nucleus="indigo")
        elif i == 1:
            add_line(slide, x + 17, 681, x + 37, 681, color="paper", width=1.8)
            add_line(slide, x + 27, 671, x + 27, 691, color="paper", width=1.8)
        else:
            add_bin_grid(slide, x + 14, 668, 7)
        add_text(slide, title, x + 50, 666, 84, 18, size=9.4, color="ink", bold=True)
        add_text(slide, body, x + 18, 704, 116, 34, size=7.3, color="muted", align=PP_ALIGN.CENTER, line_spacing=0.9)

    # Main workflow arrows
    add_arrow(slide, 560, 370, 560, 590, color="line", width=1.4)
    add_arrow(slide, 942, 372, 1238, 490, color="cyan", width=2.1)
    add_arrow(slide, 942, 595, 1238, 540, color="teal", width=1.7)

    # Outputs
    add_rect(slide, 1242, 156, 300, 590, fill="paper", line="line", radius=True, width=0.9)
    add_label_pill(slide, "Outputs and biological impact", 1260, 176, 205, fill="soft_teal", color="teal")
    out_items = [
        ("Refined whole cells", "corrected assignments with improved coherence", "cyan", 230),
        ("Reconstructed partial / anuclear cells", "profiles recovered from residual transcripts or bins", "amber", 318),
        ("Improved cell-type maps", "cleaner profile purity and reduced conflict", "indigo", 414),
        ("Downstream interpretation", "cell typing | niches | ligand-receptor | CNV / clone analysis", "teal", 512),
    ]
    for title, body, accent, y in out_items:
        add_circle(slide, 1278, y + 16, 17, fill=accent, line=accent, width=0)
        if "whole" in title:
            add_micro_cell(slide, 1278, y + 16, 0.33, fill="paper", nucleus="indigo")
        elif "partial" in title:
            add_circle(slide, 1272, y + 14, 12, fill="paper", line="paper", width=0)
            add_line(slide, 1269, y + 4, 1290, y + 29, color="paper", width=1.8)
        elif "maps" in title:
            add_bin_grid(slide, 1265, y + 3, 8)
        else:
            add_line(slide, 1264, y + 13, 1292, y + 13, color="paper", width=1.5)
            add_line(slide, 1278, y + 1, 1278, y + 29, color="paper", width=1.5)
        add_text(slide, title, 1308, y, 190, 18, size=10.0, bold=True)
        add_text(slide, body, 1308, y + 23, 176, 34, size=7.3, color="muted", line_spacing=0.9)
        if y < 512:
            add_line(slide, 1278, y + 40, 1278, y + 78, color="line2", width=1.0)

    # Bottom note
    add_rect(slide, 52, 815, 1490, 34, fill="soft_slate", line="soft_slate", radius=True, width=0)
    add_text(
        slide,
        "Core idea: combine transcript/cell spatial coherence with PMI/NPMI gene-pair evidence to prune conflicts, rescue residual transcripts, stitch compatible fragments, and recover profiles where segmentation is absent or incomplete.",
        70,
        824,
        1454,
        15,
        size=8.2,
        color="slate",
        align=PP_ALIGN.CENTER,
    )

    prs.save(PPTX_OUT)


def svg_text(dwg, text, x, y, size=16, fill="ink", weight="normal", anchor="start", linesp=1.0):
    lines = text.split("\n")
    for i, line in enumerate(lines):
        dwg.add(dwg.text(
            line,
            insert=(x, y + i * size * 1.25 * linesp),
            fill="#" + COL[fill],
            font_size=size,
            font_family="Arial, Helvetica, sans-serif",
            font_weight=weight,
            text_anchor=anchor,
        ))


def svg_round(dwg, x, y, w, h, fill="paper", stroke="line", sw=1.0, r=14):
    dwg.add(dwg.rect(insert=(x, y), size=(w, h), rx=r, ry=r, fill="#" + COL[fill], stroke="#" + COL[stroke], stroke_width=sw))


def svg_circle(dwg, cx, cy, r, fill, stroke="paper", sw=1.0):
    dwg.add(dwg.circle(center=(cx, cy), r=r, fill="#" + COL[fill], stroke="#" + COL[stroke], stroke_width=sw))


def svg_arrow(dwg, x1, y1, x2, y2, color="cyan", sw=2.0):
    dwg.add(dwg.line(start=(x1, y1), end=(x2, y2), stroke="#" + COL[color], stroke_width=sw, stroke_linecap="round"))
    ang = math.atan2(y2 - y1, x2 - x1)
    size = 13
    pts = [
        (x2, y2),
        (x2 - size * math.cos(ang - 0.46), y2 - size * math.sin(ang - 0.46)),
        (x2 - size * math.cos(ang + 0.46), y2 - size * math.sin(ang + 0.46)),
    ]
    dwg.add(dwg.polygon(points=pts, fill="#" + COL[color]))


def svg_micro_cell(dwg, cx, cy, scale, fill="soft_cyan", nucleus="indigo"):
    svg_circle(dwg, cx, cy, 26 * scale, fill=fill, stroke="cyan", sw=0.9)
    svg_circle(dwg, cx - 4 * scale, cy - 2 * scale, 9 * scale, fill=nucleus, stroke="paper", sw=0.6)
    for dx, dy, c in [(-18, -12, "coral"), (12, -15, "teal"), (19, 8, "amber"), (-10, 15, "indigo")]:
        svg_circle(dwg, cx + dx * scale, cy + dy * scale, 3.2 * scale, fill=c, stroke=c, sw=0)


def svg_bin_grid(dwg, x, y, s):
    for i in range(4):
        for j in range(4):
            fill = ["soft_cyan", "soft_indigo", "soft_teal", "soft_amber"][(i + j) % 4]
            dwg.add(dwg.rect(insert=(x + j * s, y + i * s), size=(s - 2, s - 2), fill="#" + COL[fill], stroke="#" + COL["line2"], stroke_width=0.5))
    for px, py, c in [(1.5, 1.2, "coral"), (2.4, 2.5, "teal"), (0.7, 2.2, "indigo")]:
        svg_circle(dwg, x + px * s, y + py * s, 5, fill=c, stroke=c, sw=0)


def add_svg():
    dwg = svgwrite.Drawing(str(SVG_OUT), size=(f"{W}px", f"{H}px"), viewBox=f"0 0 {W} {H}")
    dwg.add(dwg.rect(insert=(0, 0), size=(W, H), fill="#" + COL["bg"]))
    logo_data = base64.b64encode(LOGO.read_bytes()).decode("ascii")
    dwg.add(dwg.image(href=f"data:image/png;base64,{logo_data}", insert=(44, 35), size=(142, 61)))
    svg_text(dwg, "TRACER", 206, 61, size=31, fill="ink", weight="700")
    svg_text(dwg, "Transcript Coherence-Aware Reconstruction for spatial transcriptomics", 207, 89, size=13, fill="muted")
    svg_text(dwg, "From noisy 2D transcript assignments to coherent whole-cell and partial-cell profiles", 1455, 62, size=18, fill="ink", weight="700", anchor="end")
    dwg.add(dwg.line(start=(44, 113), end=(1548, 113), stroke="#" + COL["line2"], stroke_width=1))

    def card(title, sub, x, y, w, h, accent):
        dwg.add(dwg.rect(insert=(x + 4, y + 6), size=(w, h), rx=14, ry=14, fill="#" + COL["line2"], opacity=0.55))
        svg_round(dwg, x, y, w, h, fill="paper", stroke="line", sw=0.8, r=14)
        dwg.add(dwg.rect(insert=(x, y), size=(6, h), fill="#" + COL[accent]))
        svg_text(dwg, title, x + 18, y + 34, size=18, fill="ink", weight="700")
        svg_text(dwg, sub, x + 18, y + 64, size=11.2, fill="muted", linesp=0.95)

    card("Imaging-based ST", "Segmented or transcript-level input\nXenium | Xenium5K | Atera\nCosMx | MERFISH", 52, 170, 280, 170, "cyan")
    svg_micro_cell(dwg, 100, 292, 0.78)
    svg_micro_cell(dwg, 148, 286, 0.66, fill="soft_teal", nucleus="teal")
    for i, c in enumerate(["coral", "amber", "indigo", "teal", "cyan"]):
        svg_circle(dwg, 215 + i * 14, 303 - (i % 2) * 10, 4, fill=c, stroke=c, sw=0)
    card("Sequencing-based ST", "VisiumHD, bin-level or pixel-level\nNo prior segmentation required", 52, 370, 280, 155, "teal")
    svg_bin_grid(dwg, 82, 448, 18)
    svg_text(dwg, "ambiguous\nbins / pixels", 178, 475, size=11, fill="muted", linesp=0.9)
    svg_arrow(dwg, 335, 276, 390, 276, color="cyan", sw=2.2)
    svg_arrow(dwg, 335, 448, 390, 448, color="teal", sw=2.2)

    svg_round(dwg, 405, 168, 310, 132, fill="paper", stroke="line", sw=0.8)
    svg_text(dwg, "Why segmentation breaks", 423, 193, size=17, fill="ink", weight="700")
    svg_text(dwg, "3D tissue volume projected to 2D", 423, 215, size=11, fill="muted")
    for dx, dy, col in [(0, 0, "soft_indigo"), (12, 12, "soft_cyan"), (24, 24, "soft_teal")]:
        svg_round(dwg, 430 + dx, 234 + dy, 108, 42, fill=col, stroke="line2", sw=0.6, r=10)
        svg_micro_cell(dwg, 467 + dx, 254 + dy, 0.55, fill="paper", nucleus="indigo")
        svg_micro_cell(dwg, 503 + dx, 258 + dy, 0.50, fill="paper", nucleus="coral")
    dwg.add(dwg.line(start=(555, 230), end=(611, 268), stroke="#" + COL["coral"], stroke_width=1.4))
    dwg.add(dwg.line(start=(555, 286), end=(613, 258), stroke="#" + COL["coral"], stroke_width=1.4))
    svg_text(dwg, "overlap\nmixed genes\npartial / anuclear", 603, 252, size=11, fill="muted", linesp=0.88)
    svg_arrow(dwg, 720, 245, 770, 245, color="slate", sw=1.8)

    svg_round(dwg, 760, 148, 445, 222, fill="paper", stroke="line", sw=0.9)
    svg_round(dwg, 784, 165, 150, 28, fill="soft_indigo", stroke="soft_indigo", sw=0, r=14)
    svg_text(dwg, "TRACER method core", 859, 183, size=12, fill="indigo", weight="700", anchor="middle")
    steps = [
        ("1", "Coherence\nmap", "gene-pair\nevidence", "indigo", "soft_indigo"),
        ("2", "Conflict\nprune", "incompatible\nassignments", "coral", "soft_coral"),
        ("3", "Rescue &\nregroup", "spatial\nwitnesses", "teal", "soft_teal"),
        ("4", "Stitch /\nrefine", "coherent\nfragments", "cyan", "soft_cyan"),
        ("5", "Partial-cell\nreconstruct", "anuclear /\ncut cells", "amber", "soft_amber"),
    ]
    sx = 782
    for idx, (num, title, sub, accent, fill) in enumerate(steps):
        x = sx + idx * 82
        svg_circle(dwg, x + 30, 216, 23, fill=fill, stroke=accent, sw=1.1)
        svg_text(dwg, num, x + 30, 223, size=16, fill=accent, weight="700", anchor="middle")
        svg_text(dwg, title, x + 30, 262, size=12.4, fill="ink", weight="700", anchor="middle", linesp=0.82)
        svg_text(dwg, sub, x + 30, 302, size=8.4, fill="muted", anchor="middle", linesp=0.88)
        if idx < len(steps) - 1:
            svg_arrow(dwg, x + 57, 216, x + 78, 216, color="line", sw=1.2)
    svg_round(dwg, 792, 333, 378, 22, fill="soft_slate", stroke="soft_slate", sw=0, r=11)
    svg_text(dwg, "composition score C(entity) + spatial gates + witness-supported reassignment", 981, 348, size=10, fill="slate", anchor="middle")

    svg_round(dwg, 405, 595, 558, 190, fill="paper", stroke="line", sw=0.9)
    svg_text(dwg, "Operation modes", 430, 637, size=17, fill="ink", weight="700")
    modes = [
        ("Refine in place", "preserve original whole cells\nprune conflict; reconstruct partials", "cyan", "soft_cyan"),
        ("Resegment", "start from segmentation\nallow new entities and stitching", "indigo", "soft_indigo"),
        ("Noseg", "discard labels\nreconstruct from bins / pixels", "teal", "soft_teal"),
    ]
    for i, (title, body, accent, fill) in enumerate(modes):
        x = 430 + i * 172
        svg_round(dwg, x, 653, 150, 96, fill=fill, stroke=accent, sw=0.8)
        svg_circle(dwg, x + 27, 681, 14, fill=accent, stroke=accent, sw=0)
        if i == 0:
            svg_micro_cell(dwg, x + 27, 681, 0.30, fill="paper", nucleus="indigo")
        elif i == 1:
            dwg.add(dwg.line(start=(x + 17, 681), end=(x + 37, 681), stroke="#FFFFFF", stroke_width=1.8))
            dwg.add(dwg.line(start=(x + 27, 671), end=(x + 27, 691), stroke="#FFFFFF", stroke_width=1.8))
        else:
            svg_bin_grid(dwg, x + 14, 668, 7)
        svg_text(dwg, title, x + 50, 684, size=12.5, fill="ink", weight="700")
        svg_text(dwg, body, x + 75, 715, size=9.8, fill="muted", anchor="middle", linesp=0.9)

    svg_arrow(dwg, 560, 370, 560, 590, color="line", sw=1.4)
    svg_arrow(dwg, 942, 372, 1238, 490, color="cyan", sw=2.1)
    svg_arrow(dwg, 942, 595, 1238, 540, color="teal", sw=1.7)

    svg_round(dwg, 1242, 156, 300, 590, fill="paper", stroke="line", sw=0.9)
    svg_round(dwg, 1260, 176, 205, 28, fill="soft_teal", stroke="soft_teal", sw=0, r=14)
    svg_text(dwg, "Outputs and biological impact", 1363, 194, size=12, fill="teal", weight="700", anchor="middle")
    out_items = [
        ("Refined whole cells", "corrected assignments with improved coherence", "cyan", 230),
        ("Reconstructed partial / anuclear cells", "profiles recovered from residual transcripts or bins", "amber", 318),
        ("Improved cell-type maps", "cleaner profile purity and reduced conflict", "indigo", 414),
        ("Downstream interpretation", "cell typing | niches | ligand-receptor | CNV / clone analysis", "teal", 512),
    ]
    for title, body, accent, y in out_items:
        svg_circle(dwg, 1278, y + 16, 17, fill=accent, stroke=accent, sw=0)
        if "whole" in title:
            svg_micro_cell(dwg, 1278, y + 16, 0.33, fill="paper", nucleus="indigo")
        elif "maps" in title:
            svg_bin_grid(dwg, 1265, y + 3, 8)
        else:
            dwg.add(dwg.line(start=(1264, y + 14), end=(1292, y + 14), stroke="#FFFFFF", stroke_width=1.5))
            dwg.add(dwg.line(start=(1278, y + 1), end=(1278, y + 29), stroke="#FFFFFF", stroke_width=1.5))
        svg_text(dwg, title, 1308, y + 16, size=13.2, fill="ink", weight="700")
        svg_text(dwg, body, 1308, y + 39, size=9.6, fill="muted", linesp=0.9)
        if y < 512:
            dwg.add(dwg.line(start=(1278, y + 40), end=(1278, y + 78), stroke="#" + COL["line2"], stroke_width=1))

    svg_round(dwg, 52, 815, 1490, 34, fill="soft_slate", stroke="soft_slate", sw=0, r=17)
    svg_text(
        dwg,
        "Core idea: combine transcript/cell spatial coherence with PMI/NPMI gene-pair evidence to prune conflicts, rescue residual transcripts, stitch compatible fragments, and recover profiles where segmentation is absent or incomplete.",
        800,
        837,
        size=11,
        fill="slate",
        anchor="middle",
    )
    dwg.save()


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    if not LOGO.exists():
        raise FileNotFoundError(LOGO)

    # Verify the logo can be read and preserve it as the sole raster element.
    with Image.open(LOGO) as im:
        im.verify()

    add_pptx()
    add_svg()
    cairosvg.svg2pdf(url=str(SVG_OUT), write_to=str(PDF_OUT), output_width=W, output_height=H)
    cairosvg.svg2png(url=str(SVG_OUT), write_to=str(PNG_OUT), output_width=3200, output_height=1800)

    print("Created:")
    for path in (PPTX_OUT, SVG_OUT, PDF_OUT, PNG_OUT):
        print(f"  {path}")


if __name__ == "__main__":
    main()
